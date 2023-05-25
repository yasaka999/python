from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE

# 打开 PowerPoint 文件
ppt = Presentation('../files/test.pptx')

# 选择要操作的幻灯片
slide_index = 0  # 假设要操作第一张幻灯片
slide = ppt.slides[slide_index]

# 遍历幻灯片上的形状
for shape in slide.shapes:
    if shape.shape_type == MSO_SHAPE.CHART:
        chart = shape.chart

        # 打印当前图表的数据
        plot = chart.plots[0]  # 假设只有一个绘图区域
        for series in plot.series:
            print(f"Series Name: {series.name}")
            print("Data Points:")
            for point in series.points:
                print(f"  Category: {point.category}, Value: {point.value}")
