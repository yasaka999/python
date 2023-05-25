from pptx import Presentation
from pptx.enum.shapes import ShapeType

# 打开 PowerPoint 文件
ppt = Presentation('../files/test.pptx')

# 选择要操作的幻灯片
slide_index = 4  # 假设要操作第一张幻灯片
slide = ppt.slides[slide_index]

# 遍历幻灯片上的所有对象
for shape in slide.shapes:
    if shape.shape_type == ShapeType.TEXT_BOX:
        # 打印文本框对象
        text_frame = shape.text_frame
        print(f"Text Box: {text_frame.text}")
    elif shape.shape_type == ShapeType.CHART:
        # 打印图表对象
        chart = shape.chart
        print(f"Chart: {chart.name}")
    else:
        # 其他形状类型的对象
        print(f"Shape: {shape.shape_type}")

# 保存修改后的 PowerPoint 文件
ppt.save('../files/modified2.pptx')
