import requests
import sys
from pptx import Presentation

openai_api_key = "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"
endpoint = "https://api.openai.com/v1/chat/completions"
headers = {
    "Content-Type": "application/json",
    "Authorization": f"Bearer {openai_api_key}"
}

def translate_paragraph_with_retry(english_paragraph):
    data = {
        "model": "gpt-3.5-turbo",
        "messages": [{
            "role": "user",
            "content": f"Translate the following content from English to Chinese:\n{english_paragraph}"
        }],
    }

    translated_paragraph = None
    retry_count = 0

    while retry_count < 3:  # 最多重试3次
        try:
            print (english_paragraph)
            response = requests.post(endpoint, json=data, headers=headers, timeout=15)
            response.raise_for_status()
            result = response.json()
            translated_paragraph = result['choices'][0]['message']['content']
            print (translated_paragraph)
            break  # 如果成功获取响应，跳出循环
        except requests.exceptions.RequestException as e:
            print(f"API请求失败: {str(e)}")
            retry_count += 1

    if translated_paragraph is None:
        translated_paragraph = english_paragraph  # 如果请求失败，返回原段落

    return translated_paragraph

def translate_presentation_to_chinese(input_presentation_file, output_ppt_file2):
    ppt = Presentation(input_presentation_file)

    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                for paragraph in shape.text_frame.paragraphs:
                    #print(paragraph.runs)
                    english_paragraph = "".join(run.text for run in paragraph.runs)
                    if english_paragraph == "" or english_paragraph == "\n" or english_paragraph.strip() == "":
                        translated_paragraph = english_paragraph
                    else:
                        translated_paragraph = translate_paragraph_with_retry(english_paragraph)
                    # 分割翻译后的段落，以匹配运行
                    #translated_lines = translated_paragraph.split('\n')
                    #print (translated_lines)
                    

                    # 更新每个运行的文本
                    #for run, translated_line in zip(paragraph.runs, translated_lines):
                        #run.text = translated_paragraph
                        #print(run.text)
                    if len(paragraph.runs) > 0:
                        first_run = paragraph.runs[0]  # 选择第一个运行
                        first_run.text = translated_paragraph
                        for run in paragraph.runs[1:]:
                            run.text = ''    

    translated_pptx_file = sys.argv[2]
    ppt.save(translated_pptx_file)

if __name__ == '__main__':
    if len(sys.argv) != 3:
        print("Usage: python program_fixer.py input_file output_file")
        sys.exit(1)

    translate_presentation_to_chinese(sys.argv[1], sys.argv[2])
