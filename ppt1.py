from pptx import Presentation
import datetime
from pptx.dml.color import RGBColor

# 获取当前日期
current_date = datetime.date.today()

# 计算上一周的起始日期和结束日期
start_of_last_week = current_date - datetime.timedelta(days=current_date.weekday() + 7)
end_of_last_week = start_of_last_week + datetime.timedelta(days=6)

# 格式化日期为字符串
start_of_last_week_str = start_of_last_week.strftime("%Y年%m月%d日")
end_of_last_week_str = end_of_last_week.strftime("%Y年%m月%d日")

# 生成周期时间字符串
period_str = f"{start_of_last_week_str}~{end_of_last_week_str}"

print(period_str)


# 打开 PowerPoint 文件
ppt = Presentation('../files/test.pptx')

# 选择要操作的幻灯片
slide_index = 0  # 假设要操作第一张幻灯片
slide = ppt.slides[slide_index]

# 遍历幻灯片上的形状
for shape in slide.shapes:
    if shape.has_text_frame:
        text_frame = shape.text_frame

        if "~" in text_frame.text:
            # 获取原始文本框样式
            original_font = text_frame.text_frame
            original_text = text_frame.text

            # 替换文本框内容
            text_frame.clear()  # 清空文本框内容
            new_paragraph = text_frame.add_paragraph()
            new_paragraph.text = period_str

            # 设置保持原样式的样式属性
            text_frame.text_frame = original_font
            text_frame.text = original_text



# 保存修改后的 PowerPoint 文件
ppt.save('../files/modified.pptx')
